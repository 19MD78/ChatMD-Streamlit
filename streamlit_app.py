from __future__ import annotations

import base64
import hashlib
import io
import datetime as dt
import hmac
import html
import json
import sqlite3
import time
import uuid
from typing import Any, Dict, List, Optional, Tuple

import streamlit as st

try:
    from openai import OpenAI
except Exception:
    OpenAI = None

try:
    from google import genai
    from google.genai import types as genai_types
except Exception:
    genai = None
    genai_types = None

try:
    from PyPDF2 import PdfReader
except Exception:
    PdfReader = None

try:
    import docx
except Exception:
    docx = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from streamlit_paste_button import paste_image_button
except Exception:
    paste_image_button = None


APP_NAME = "ChatMD"
APP_VERSION = "V. 2026_05_25_6"
DEFAULT_PROVIDER = "Google Gemini"
HISTORY_DB_PATH = "chatmd_history.db"

FALLBACK_MODELS: Dict[str, List[str]] = {
    "Google Gemini": [
        "gemini-3-pro-preview",
        "gemini-2.5-pro",
        "gemini-2.5-flash",
        "gemini-2.0-flash",
    ],
    "OpenAI": [
        "gpt-4o",
        "gpt-4o-mini",
        "gpt-4.1",
        "gpt-4.1-mini",
    ],
}

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".gif"}
IMAGE_MIME_TYPES = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
    ".gif": "image/gif",
}
MAX_INLINE_IMAGE_BYTES = 20 * 1024 * 1024
TEXT_FILE_CHAR_LIMIT = 1_000_000


DEEP_RESEARCH_AGENT_FAST = "deep-research-preview-04-2026"
DEEP_RESEARCH_POLL_SECONDS = 10
DEEP_RESEARCH_MAX_POLLS = 30

DEFAULT_AGENTS = {
    "Bendras asistentas": {
        "icon": "💬",
        "system_prompt": "Tu esi naudingas, tikslus ir aiškiai lietuviškai atsakantis asistentas.",
        "examples": "",
    },
    "Darbo pagalbininkas": {
        "icon": "🧩",
        "system_prompt": "Padėk struktūruoti užduotis, rašyti aiškiai, trumpai ir profesionaliai lietuviškai.",
        "examples": "Atsakymo struktūra: trumpa išvada, tada konkretūs veiksmai, tada rizikos arba pastabos.",
    },
}


def get_secret(name: str) -> str:
    try:
        return str(st.secrets.get(name, "") or "").strip()
    except Exception:
        return ""


def get_api_key(provider: str) -> str:
    if provider == "Google Gemini":
        return get_secret("GOOGLE_API_KEY")
    if provider == "OpenAI":
        return get_secret("OPENAI_API_KEY")
    return ""


def now_iso() -> str:
    return dt.datetime.now().isoformat(timespec="seconds")


def init_history_db() -> None:
    conn = sqlite3.connect(HISTORY_DB_PATH)
    try:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS chats (
                id TEXT PRIMARY KEY,
                title TEXT NOT NULL,
                messages_json TEXT NOT NULL,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL,
                mode TEXT NOT NULL,
                provider TEXT NOT NULL,
                model TEXT NOT NULL,
                agent_name TEXT NOT NULL DEFAULT ''
            )
            """
        )

        existing_columns = {
            row[1] for row in conn.execute("PRAGMA table_info(chats)").fetchall()
        }

        if "agent_name" not in existing_columns:
            conn.execute("ALTER TABLE chats ADD COLUMN agent_name TEXT NOT NULL DEFAULT ''")

        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS agents (
                name TEXT PRIMARY KEY,
                icon TEXT NOT NULL,
                system_prompt TEXT NOT NULL,
                examples TEXT NOT NULL,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL
            )
            """
        )

        existing_agents = conn.execute("SELECT COUNT(*) FROM agents").fetchone()[0]

        if existing_agents == 0:
            created_at = now_iso()
            for name, agent in DEFAULT_AGENTS.items():
                conn.execute(
                    """
                    INSERT INTO agents (name, icon, system_prompt, examples, created_at, updated_at)
                    VALUES (?, ?, ?, ?, ?, ?)
                    """,
                    (
                        name,
                        agent.get("icon", "🤖"),
                        agent.get("system_prompt", ""),
                        agent.get("examples", ""),
                        created_at,
                        created_at,
                    ),
                )

        conn.commit()
    finally:
        conn.close()


def make_chat_title(messages: List[Dict[str, str]]) -> str:
    for message in messages:
        if message.get("role") == "user":
            title = str(message.get("content", "")).strip().replace("\n", " ")
            return title[:60] if title else "Pokalbis"
    return "Pokalbis"


def save_current_chat_to_db() -> None:
    messages = st.session_state.get("messages", [])

    if not messages:
        return

    init_history_db()

    chat_id = st.session_state.get("current_chat_id") or str(uuid.uuid4())
    st.session_state.current_chat_id = chat_id

    title = make_chat_title(messages)
    updated_at = now_iso()
    messages_json = json.dumps(messages, ensure_ascii=False)

    conn = sqlite3.connect(HISTORY_DB_PATH)
    try:
        existing = conn.execute(
            "SELECT created_at FROM chats WHERE id = ?",
            (chat_id,),
        ).fetchone()

        created_at = existing[0] if existing else updated_at

        conn.execute(
            """
            INSERT INTO chats (
                id, title, messages_json, created_at, updated_at, mode, provider, model, agent_name
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            ON CONFLICT(id) DO UPDATE SET
                title = excluded.title,
                messages_json = excluded.messages_json,
                updated_at = excluded.updated_at,
                mode = excluded.mode,
                provider = excluded.provider,
                model = excluded.model,
                agent_name = excluded.agent_name
            """,
            (
                chat_id,
                title,
                messages_json,
                created_at,
                updated_at,
                st.session_state.get("mode", "chat"),
                st.session_state.get("provider", DEFAULT_PROVIDER),
                st.session_state.get("model", ""),
                st.session_state.get("agent_name", ""),
            ),
        )
        conn.commit()
    finally:
        conn.close()


def list_saved_chats(limit: int = 30) -> List[Dict[str, str]]:
    init_history_db()

    conn = sqlite3.connect(HISTORY_DB_PATH)
    conn.row_factory = sqlite3.Row
    try:
        rows = conn.execute(
            """
            SELECT id, title, created_at, updated_at, mode, provider, model, agent_name
            FROM chats
            ORDER BY updated_at DESC
            LIMIT ?
            """,
            (limit,),
        ).fetchall()

        return [dict(row) for row in rows]
    finally:
        conn.close()


def load_chat_from_db(chat_id: str) -> Optional[Dict[str, Any]]:
    init_history_db()

    conn = sqlite3.connect(HISTORY_DB_PATH)
    conn.row_factory = sqlite3.Row
    try:
        row = conn.execute(
            "SELECT * FROM chats WHERE id = ?",
            (chat_id,),
        ).fetchone()

        if not row:
            return None

        chat = dict(row)
        chat["messages"] = json.loads(chat.get("messages_json", "[]"))
        return chat
    finally:
        conn.close()


def delete_chat_from_db(chat_id: str) -> None:
    init_history_db()

    conn = sqlite3.connect(HISTORY_DB_PATH)
    try:
        conn.execute("DELETE FROM chats WHERE id = ?", (chat_id,))
        conn.commit()
    finally:
        conn.close()


def load_agents_from_db() -> Dict[str, Dict[str, str]]:
    init_history_db()

    conn = sqlite3.connect(HISTORY_DB_PATH)
    conn.row_factory = sqlite3.Row
    try:
        rows = conn.execute(
            """
            SELECT name, icon, system_prompt, examples
            FROM agents
            ORDER BY created_at ASC, name ASC
            """
        ).fetchall()

        agents = {
            row["name"]: {
                "icon": row["icon"],
                "system_prompt": row["system_prompt"],
                "examples": row["examples"],
            }
            for row in rows
        }

        return agents or {name: dict(agent) for name, agent in DEFAULT_AGENTS.items()}
    finally:
        conn.close()


def save_agent_to_db(name: str, icon: str, system_prompt: str, examples: str, old_name: str = "") -> None:
    init_history_db()

    name = name.strip()
    icon = icon.strip() or "🤖"
    system_prompt = system_prompt.strip()
    examples = examples.strip()
    updated_at = now_iso()

    if not name:
        raise ValueError("Agento pavadinimas negali būti tuščias.")

    if not system_prompt:
        raise ValueError("Agento taisyklės negali būti tuščios.")

    conn = sqlite3.connect(HISTORY_DB_PATH)
    try:
        if old_name and old_name != name:
            existing = conn.execute("SELECT 1 FROM agents WHERE name = ?", (name,)).fetchone()
            if existing:
                raise ValueError("Agentas tokiu pavadinimu jau yra.")

            old_row = conn.execute("SELECT created_at FROM agents WHERE name = ?", (old_name,)).fetchone()
            created_at = old_row[0] if old_row else updated_at
            conn.execute("DELETE FROM agents WHERE name = ?", (old_name,))
        else:
            row = conn.execute("SELECT created_at FROM agents WHERE name = ?", (name,)).fetchone()
            created_at = row[0] if row else updated_at

        conn.execute(
            """
            INSERT INTO agents (name, icon, system_prompt, examples, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?, ?)
            ON CONFLICT(name) DO UPDATE SET
                icon = excluded.icon,
                system_prompt = excluded.system_prompt,
                examples = excluded.examples,
                updated_at = excluded.updated_at
            """,
            (name, icon, system_prompt, examples, created_at, updated_at),
        )
        conn.commit()
    finally:
        conn.close()


def delete_agent_from_db(name: str) -> None:
    init_history_db()

    conn = sqlite3.connect(HISTORY_DB_PATH)
    try:
        conn.execute("DELETE FROM agents WHERE name = ?", (name,))
        conn.commit()
    finally:
        conn.close()


def make_agents_export_json() -> str:
    agents = load_agents_from_db()

    payload = {
        "app": APP_NAME,
        "version": APP_VERSION,
        "exported_at": now_iso(),
        "agents": [
            {
                "name": name,
                "icon": agent.get("icon", "🤖"),
                "system_prompt": agent.get("system_prompt", ""),
                "examples": agent.get("examples", ""),
            }
            for name, agent in agents.items()
        ],
    }

    return json.dumps(payload, ensure_ascii=False, indent=2)


def import_agents_from_json_bytes(data: bytes) -> int:
    try:
        text = data.decode("utf-8")
    except Exception:
        text = data.decode("utf-8-sig", errors="ignore")

    payload = json.loads(text)

    if isinstance(payload, dict):
        agents_data = payload.get("agents", [])
    elif isinstance(payload, list):
        agents_data = payload
    else:
        raise ValueError("JSON formatas netinkamas.")

    if not isinstance(agents_data, list):
        raise ValueError("JSON faile nerastas agentų sąrašas.")

    imported_count = 0

    for item in agents_data:
        if not isinstance(item, dict):
            continue

        name = str(item.get("name", "") or "").strip()
        icon = str(item.get("icon", "🤖") or "🤖").strip()
        system_prompt = str(item.get("system_prompt", "") or "").strip()
        examples = str(item.get("examples", "") or "").strip()

        if not name or not system_prompt:
            continue

        save_agent_to_db(
            name=name,
            icon=icon,
            system_prompt=system_prompt,
            examples=examples,
        )
        imported_count += 1

    if imported_count == 0:
        raise ValueError("Nepavyko importuoti nė vieno agento. Patikrink JSON struktūrą.")

    return imported_count


def render_css() -> None:
    st.markdown(
        """
        <style>
        html, body, [class*="css"] {
            font-size: 14px;
        }

        .stApp {
            background: #f7f9fc;
            color: #111827;
        }

        .main .block-container {
            max-width: 980px;
            padding-top: 1.0rem;
            padding-bottom: 7rem;
        }

        [data-testid="stSidebar"] {
            background: #eef3fa;
            border-right: 1px solid #d8e2f0;
            width: 250px !important;
            min-width: 250px !important;
            max-width: 250px !important;
        }

        [data-testid="stSidebar"] > div:first-child {
            width: 250px !important;
            min-width: 250px !important;
            max-width: 250px !important;
        }

        [data-testid="stSidebar"] * {
            font-size: 0.90rem;
        }

        [data-testid="stSidebar"] hr {
            margin: 0.65rem 0 !important;
        }

        [data-testid="stSidebar"] h1,
        [data-testid="stSidebar"] h2,
        [data-testid="stSidebar"] h3 {
            margin-top: 0.35rem !important;
            margin-bottom: 0.35rem !important;
        }

        [data-testid="stSidebar"] .stMarkdown {
            margin-bottom: 0.15rem !important;
        }

        [data-testid="stSidebar"] button p,
        [data-testid="stSidebar"] label p,
        [data-testid="stSidebar"] .stMarkdown p {
            white-space: nowrap !important;
            overflow: hidden !important;
            text-overflow: ellipsis !important;
        }

        [data-testid="stSidebar"] .stButton > button {
            background: #e8f1ff !important;
            border: 1px solid #cfe0ff !important;
            border-radius: 13px !important;
            color: #0f172a !important;
            box-shadow: none !important;
            min-height: 2.15rem !important;
            transition: transform 0.08s ease, background-color 0.08s ease, border-color 0.08s ease, box-shadow 0.08s ease !important;
        }

        [data-testid="stSidebar"] .stButton > button:hover {
            background: #dbeafe !important;
            border-color: #bfdbfe !important;
            color: #0f172a !important;
            box-shadow: 0 3px 10px rgba(37, 99, 235, 0.08) !important;
        }

        [data-testid="stSidebar"] .stButton > button:active {
            transform: scale(0.97) !important;
            background: #bfdbfe !important;
            border-color: #93c5fd !important;
            box-shadow: inset 0 0 0 1px #93c5fd !important;
        }

        [data-testid="stSidebar"] .stButton > button:focus {
            outline: none !important;
            box-shadow: 0 0 0 2px rgba(37, 99, 235, 0.18) !important;
        }

        [data-testid="stSidebar"] .stButton > button p {
            color: #0f172a !important;
            font-weight: 500 !important;
        }

        h1 {
            font-size: 1.55rem !important;
            margin-bottom: 0.25rem !important;
        }

        h2, h3 {
            font-size: 1.05rem !important;
        }

        .chat-title {
            display: flex;
            align-items: center;
            gap: 0.48rem;
            padding: 0 0 0.12rem 0;
            margin-top: -0.35rem;
        }

        .chat-logo {
            width: 34px;
            height: 34px;
            object-fit: contain;
        }

        .brand-main {
            font-size: 1.14rem;
            font-weight: 700;
            line-height: 1;
        }

        .brand-sub {
            font-size: 0.68rem;
            color: #526173;
            margin-top: 0.05rem;
        }

        .version-pill {
            display: inline-block;
            padding: 0.13rem 0.42rem;
            border-radius: 999px;
            background: #fff7ed;
            color: #9a3412;
            border: 1px solid #fed7aa;
            font-size: 0.70rem;
            margin-bottom: 0.35rem;
        }

        .status-pill {
            display: inline-block;
            padding: 0.17rem 0.5rem;
            border-radius: 999px;
            background: #dbeafe;
            color: #1d4ed8;
            font-size: 0.76rem;
            margin-right: 0.28rem;
            margin-bottom: 0.2rem;
        }

        [data-testid="stSidebar"] .stSelectbox,
        [data-testid="stSidebar"] .stSlider,
        [data-testid="stSidebar"] .stToggle,
        [data-testid="stSidebar"] .stRadio {
            margin-bottom: 0.25rem !important;
        }

        [data-testid="stSidebar"] label {
            margin-bottom: 0.15rem !important;
        }

        .chat-area {
            margin-top: 1.1rem;
            padding-bottom: 1.5rem;
        }

        .msg-row {
            display: flex;
            width: 100%;
            margin: 0.45rem 0;
        }

        .msg-row.user {
            justify-content: flex-end;
        }

        .msg-bubble {
            max-width: 78%;
            padding: 0.62rem 0.78rem;
            border-radius: 16px;
            font-size: 0.96rem;
            line-height: 1.46;
            white-space: pre-wrap;
            word-wrap: break-word;
        }

        .msg-bubble.user {
            background: #e8f1ff;
            border: 1px solid #cfe0ff;
            color: #0f172a;
            border-bottom-right-radius: 6px;
        }

        .thinking-wrap {
            display: flex;
            justify-content: flex-start;
            width: 100%;
            margin: 0.45rem 0;
        }

        .thinking-card {
            display: inline-flex;
            align-items: center;
            gap: 0.65rem;
            padding: 0.62rem 0.78rem;
            border-radius: 16px;
            background: #ffffff;
            border: 1px solid #d8e2f0;
            box-shadow: 0 5px 18px rgba(15, 23, 42, 0.05);
            color: #334155;
            font-size: 0.92rem;
        }

        .thinking-spinner {
            width: 16px;
            height: 16px;
            border: 2px solid #cbd5e1;
            border-top-color: #2563eb;
            border-radius: 50%;
            animation: chatmd-spin 0.8s linear infinite;
        }

        .thinking-text {
            display: inline-flex;
            align-items: center;
            gap: 0.15rem;
        }

        .thinking-dots span {
            animation: chatmd-blink 1.2s infinite;
            opacity: 0.25;
        }

        .thinking-dots span:nth-child(2) {
            animation-delay: 0.2s;
        }

        .thinking-dots span:nth-child(3) {
            animation-delay: 0.4s;
        }

        .thinking-time {
            color: #64748b;
            font-size: 0.78rem;
            margin-top: 0.05rem;
        }

        @keyframes chatmd-spin {
            to {
                transform: rotate(360deg);
            }
        }

        @keyframes chatmd-blink {
            0%, 80%, 100% {
                opacity: 0.25;
            }
            40% {
                opacity: 1;
            }
        }

        [data-testid="stChatInput"] {
            max-width: 980px;
            margin: 0 auto;
        }

        [data-testid="stChatInput"] > div {
            background: #eef4fb !important;
            border: 1px solid #cbdcf0 !important;
            border-radius: 14px !important;
        }

        [data-testid="stChatInput"] button {
            border-radius: 11px !important;
            opacity: 1 !important;
            border: 0 !important;
            outline: 0 !important;
            box-shadow: none !important;
            transition: transform 0.08s ease, background-color 0.08s ease, box-shadow 0.08s ease !important;
        }

        [data-testid="stChatInput"] button svg {
            opacity: 1 !important;
        }

        [data-testid="stChatInput"] button:first-child,
        [data-testid="stChatInput"] button:first-child:hover,
        [data-testid="stChatInput"] button:first-child:focus {
            background: #dbeafe !important;
            border: 0 !important;
            border-right: 0 !important;
            box-shadow: none !important;
            margin-right: 0.35rem !important;
            color: #2563eb !important;
        }

        [data-testid="stChatInput"] button:last-child,
        [data-testid="stChatInput"] button:last-child:hover,
        [data-testid="stChatInput"] button:last-child:focus {
            background: #dbeafe !important;
            border: 0 !important;
            box-shadow: none !important;
            color: #1d4ed8 !important;
            opacity: 1 !important;
            filter: none !important;
        }

        [data-testid="stChatInput"] button:active {
            transform: scale(0.92) !important;
            background: #bfdbfe !important;
            box-shadow: inset 0 0 0 1px #93c5fd !important;
        }

        [data-testid="stChatInput"] button:last-child svg,
        [data-testid="stChatInput"] button:last-child:hover svg,
        [data-testid="stChatInput"] button:last-child:focus svg,
        [data-testid="stChatInput"] button:last-child:active svg {
            color: #1d4ed8 !important;
            stroke: #1d4ed8 !important;
            opacity: 1 !important;
            filter: none !important;
        }

        .login-wrap {
            display: flex;
            justify-content: center;
            margin-top: 12vh;
        }

        .login-card {
            width: 100%;
            max-width: 420px;
            background: white;
            border: 1px solid #d8e2f0;
            border-radius: 20px;
            padding: 1.2rem;
            box-shadow: 0 10px 30px rgba(15, 23, 42, 0.08);
            text-align: center;
        }

        .login-title {
            font-size: 1.55rem;
            font-weight: 750;
            margin-bottom: 0.35rem;
        }

        .login-subtitle {
            font-size: 0.9rem;
            color: #64748b;
        }

        @media (max-width: 720px) {
            .main .block-container {
                padding-left: 0.75rem;
                padding-right: 0.75rem;
            }

            .msg-bubble {
                max-width: 92%;
                font-size: 0.95rem;
            }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def check_password() -> bool:
    app_password = get_secret("APP_PASSWORD")

    if not app_password:
        st.error(
            "APP_PASSWORD nėra nustatytas Streamlit Secrets. "
            "Programa saugumo sumetimais nerodoma."
        )
        return False

    if st.session_state.get("authenticated", False):
        return True

    render_css()

    st.markdown(
        """
        <div class="login-wrap">
            <div class="login-card">
                <div class="login-title">🔐 ChatMD</div>
                <div class="login-subtitle">Įvesk slaptažodį, kad galėtum naudotis programa.</div>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    left, center, right = st.columns([1, 0.42, 1])

    with center:
        entered_password = st.text_input(
            "Slaptažodis",
            type="password",
            placeholder="Slaptažodis",
            label_visibility="collapsed",
        )

        if st.button("Prisijungti", use_container_width=True):
            if hmac.compare_digest(entered_password, app_password):
                st.session_state.authenticated = True
                st.rerun()
            else:
                st.error("Neteisingas slaptažodis.")

    return False


def init_state() -> None:
    st.session_state.setdefault("provider", DEFAULT_PROVIDER)
    st.session_state.setdefault("models_cache", {})
    st.session_state.setdefault("model", FALLBACK_MODELS[DEFAULT_PROVIDER][0])
    st.session_state.setdefault("web_search_enabled", True)
    st.session_state.setdefault("temperature", 0.7)
    st.session_state.setdefault("mode", "chat")
    st.session_state.setdefault("agent_name", "Bendras asistentas")
    st.session_state.setdefault("agents", load_agents_from_db())
    st.session_state.setdefault("messages", [])
    st.session_state.setdefault("chat_counter", 1)
    st.session_state.setdefault("history", [])
    st.session_state.setdefault("current_chat_id", None)
    st.session_state.setdefault("authenticated", False)
    st.session_state.setdefault("last_chat_settings_loaded", False)
    st.session_state.setdefault("pasted_images", [])


def reset_chat() -> None:
    save_current_chat_to_db()
    st.session_state.messages = []
    st.session_state.current_chat_id = None
    st.session_state.chat_counter += 1


def runtime_context() -> str:
    today = dt.datetime.now().strftime("%Y-%m-%d")
    return (
        f"Šiandienos data: {today}. "
        "Jei klausimas apie dabartinius įvykius ar faktus, naudok interneto paiešką, kai ji įjungta. "
        "Jei paieška nepavyksta, aiškiai tai pasakyk."
    )


def list_models(provider: str) -> Tuple[List[str], str]:
    api_key = get_api_key(provider)

    if not api_key:
        return FALLBACK_MODELS.get(provider, []), "fallback: nėra API rakto"

    if provider == "Google Gemini":
        if genai is None:
            return FALLBACK_MODELS[provider], "fallback: neįdiegta google-genai"

        try:
            client = genai.Client(api_key=api_key)
            names = []

            for model in client.models.list():
                name = getattr(model, "name", "") or ""

                if name.startswith("models/"):
                    name = name.split("/", 1)[1]

                if name and "embedding" not in name.lower():
                    names.append(name)

            return sorted(set(names)) or FALLBACK_MODELS[provider], "live"
        except Exception as exc:
            return FALLBACK_MODELS[provider], f"fallback: {exc}"

    if provider == "OpenAI":
        if OpenAI is None:
            return FALLBACK_MODELS[provider], "fallback: neįdiegta openai"

        try:
            client = OpenAI(api_key=api_key)
            models = client.models.list()
            names = sorted({m.id for m in models.data if getattr(m, "id", None)})
            names = [m for m in names if m.startswith(("gpt-", "o", "chatgpt-"))]
            return names or FALLBACK_MODELS[provider], "live"
        except Exception as exc:
            return FALLBACK_MODELS[provider], f"fallback: {exc}"

    return [], "fallback"


class PathLike:
    def __init__(self, name: str) -> None:
        self.name = name

    @property
    def suffix(self) -> str:
        if "." not in self.name:
            return ""
        return "." + self.name.rsplit(".", 1)[1].lower()


class PastedImageFile:
    def __init__(self, data: bytes, name: str) -> None:
        self.name = name
        self._data = data

    def getvalue(self) -> bytes:
        return self._data


def image_to_png_bytes(image) -> bytes:
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def add_pasted_image_to_state(image) -> bool:
    data = image_to_png_bytes(image)
    digest = hashlib.sha256(data).hexdigest()
    pasted_images = list(st.session_state.get("pasted_images", []))

    for item in pasted_images:
        if item.get("digest") == digest:
            return False

    name = f"clipboard_screenshot_{len(pasted_images) + 1}.png"
    pasted_images.append({"name": name, "data": data, "digest": digest})
    st.session_state.pasted_images = pasted_images
    return True


def get_pasted_image_files() -> List[PastedImageFile]:
    files = []

    for item in st.session_state.get("pasted_images", []):
        name = str(item.get("name", "clipboard_screenshot.png"))
        data = item.get("data", b"")

        if data:
            files.append(PastedImageFile(data=data, name=name))

    return files


def clear_pasted_images() -> None:
    st.session_state.pasted_images = []


def is_image_file(uploaded_file) -> bool:
    return PathLike(uploaded_file.name).suffix in IMAGE_SUFFIXES


def get_image_mime_type(uploaded_file) -> str:
    suffix = PathLike(uploaded_file.name).suffix
    return IMAGE_MIME_TYPES.get(suffix, "application/octet-stream")


def make_image_note(uploaded_files) -> str:
    if not uploaded_files:
        return ""

    image_names = [f.name for f in uploaded_files if is_image_file(f)]

    if not image_names:
        return ""

    return "Prisegti vaizdai, perduoti modeliui kaip tikri multimodal vaizdo duomenys:\n" + "\n".join(
        f"- {name}" for name in image_names
    )


def make_google_contents(messages: List[Dict[str, str]], files_context: str, uploaded_files) -> List[Any]:
    text_parts = []

    for m in messages[-20:]:
        role = "Vartotojas" if m["role"] == "user" else "Asistentas"
        text_parts.append(f"{role}: {m['content']}")

    if files_context:
        text_parts.append("Prisegtų failų tekstinis turinys:\n" + files_context)

    image_note = make_image_note(uploaded_files)

    if image_note:
        text_parts.append(image_note)

    contents: List[Any] = ["\n\n".join(text_parts)]

    if genai_types is not None and uploaded_files:
        for uploaded_file in uploaded_files:
            if not is_image_file(uploaded_file):
                continue

            image_bytes = uploaded_file.getvalue()

            if len(image_bytes) > MAX_INLINE_IMAGE_BYTES:
                contents.append(
                    f"[Vaizdas {uploaded_file.name} per didelis inline siuntimui: "
                    f"{len(image_bytes) / (1024 * 1024):.1f} MB. Didžiausia riba: 20 MB.]"
                )
                continue

            contents.append(
                genai_types.Part.from_bytes(
                    data=image_bytes,
                    mime_type=get_image_mime_type(uploaded_file),
                )
            )

    return contents


def make_openai_input_messages(
    messages: List[Dict[str, str]],
    files_context: str,
    uploaded_files,
) -> List[Dict[str, Any]]:
    input_messages: List[Dict[str, Any]] = [
        {"role": "system", "content": build_system_prompt()}
    ]
    input_messages.extend(messages[-20:])

    content_parts: List[Dict[str, str]] = []

    if files_context:
        content_parts.append(
            {
                "type": "input_text",
                "text": "Prisegtų failų tekstinis turinys:\n" + files_context,
            }
        )

    image_note = make_image_note(uploaded_files)

    if image_note:
        content_parts.append({"type": "input_text", "text": image_note})

    if uploaded_files:
        for uploaded_file in uploaded_files:
            if not is_image_file(uploaded_file):
                continue

            image_bytes = uploaded_file.getvalue()

            if len(image_bytes) > MAX_INLINE_IMAGE_BYTES:
                content_parts.append(
                    {
                        "type": "input_text",
                        "text": (
                            f"[Vaizdas {uploaded_file.name} per didelis inline siuntimui: "
                            f"{len(image_bytes) / (1024 * 1024):.1f} MB. Didžiausia riba: 20 MB.]"
                        ),
                    }
                )
                continue

            encoded = base64.b64encode(image_bytes).decode("utf-8")
            mime_type = get_image_mime_type(uploaded_file)
            data_url = f"data:{mime_type};base64,{encoded}"
            content_parts.append({"type": "input_image", "image_url": data_url})

    if content_parts:
        input_messages.append({"role": "user", "content": content_parts})

    return input_messages


def limit_text_for_model(name: str, text: str, limit: int = TEXT_FILE_CHAR_LIMIT) -> str:
    if len(text) <= limit:
        return f"[{name}]\n{text}"

    half = max(1, limit // 2)
    omitted_chars = len(text) - limit

    return (
        f"[{name}]\n"
        f"[Pastaba: failas buvo per ilgas. Modeliui perduota pradžia ir pabaiga, "
        f"praleista apie {omitted_chars:,} simbolių.]\n\n"
        f"--- FAILO PRADŽIA ---\n"
        f"{text[:half]}\n\n"
        f"--- FAILO VIDURYS PRALEISTAS ---\n\n"
        f"--- FAILO PABAIGA ---\n"
        f"{text[-half:]}"
    )


def read_uploaded_file(uploaded_file) -> str:
    name = uploaded_file.name
    suffix = PathLike(name).suffix
    data = uploaded_file.getvalue()

    try:
        if suffix in IMAGE_SUFFIXES:
            return f"[Vaizdas {name}] Vaizdas bus perduotas modeliui kaip multimodal image input, ne kaip tekstinė base64 ištrauka."

        if suffix == ".pdf":
            if PdfReader is None:
                return f"[{name}] PDF nuskaitymui reikia PyPDF2."

            import io

            reader = PdfReader(io.BytesIO(data))
            pages = [(p.extract_text() or "") for p in reader.pages]
            return limit_text_for_model(name, "\n".join(pages))

        if suffix == ".docx":
            if docx is None:
                return f"[{name}] DOCX nuskaitymui reikia python-docx."

            import io

            document = docx.Document(io.BytesIO(data))
            text = "\n".join(p.text for p in document.paragraphs if p.text.strip())
            return limit_text_for_model(name, text)

        if suffix == ".xlsx":
            if openpyxl is None:
                return f"[{name}] XLSX nuskaitymui reikia openpyxl."

            import io

            wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
            rows = []

            for ws in wb.worksheets:
                rows.append(f"Lapas: {ws.title}")

                for row in ws.iter_rows(values_only=True):
                    vals = [str(v) for v in row if v is not None]

                    if vals:
                        rows.append(" | ".join(vals))

            return limit_text_for_model(name, "\n".join(rows))

        if suffix == ".pptx":
            if Presentation is None:
                return f"[{name}] PPTX nuskaitymui reikia python-pptx."

            import io

            prs = Presentation(io.BytesIO(data))
            slides = []

            for i, slide in enumerate(prs.slides, 1):
                texts = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)

                if texts:
                    slides.append(f"Skaidrė {i}:\n" + "\n".join(texts))

            return limit_text_for_model(name, "\n\n".join(slides))

        text = data.decode("utf-8", errors="ignore")
        return limit_text_for_model(name, text)
    except Exception as exc:
        return f"[{name}] Nepavyko nuskaityti: {exc}"


def build_system_prompt() -> str:
    base = runtime_context()

    if st.session_state.mode == "agent":
        agents = st.session_state.get("agents") or load_agents_from_db()
        agent = agents.get(
            st.session_state.agent_name,
            DEFAULT_AGENTS["Bendras asistentas"],
        )
        prompt = agent.get("system_prompt", "")
        examples = agent.get("examples", "")

        if examples.strip():
            prompt += "\n\nAgento pavyzdžiai / žinios:\n" + examples.strip()

        return base + "\n\nTu esi ChatMD agentų režime.\n\n" + prompt

    if st.session_state.mode == "thinking":
        return (
            base
            + "\n\nTu esi ChatMD Mąstymo režime. Atsakyk giliau nei paprastame pokalbyje: "
            "pirmiausia įvertink klausimo esmę, tada pateik struktūruotą, aiškų ir praktišką atsakymą. "
            "Jei reikia, naudok interneto paiešką, kai ji įjungta. Neužtęsk be reikalo, bet parodyk svarbiausią logiką ir išvadas. "
            "Atsakyk lietuviškai."
        )

    if st.session_state.mode == "deep":
        return (
            base
            + "\n\nTu esi ChatMD Gilios analizės režime. Atlik išsamesnį tyrimą lietuviškai: "
            "pateik aiškią struktūrą, svarbiausias išvadas, argumentus ir šaltinius, jei jie prieinami."
        )

    return base + "\n\nTu esi naudingas, tikslus ir aiškiai lietuviškai atsakantis asistentas."


def make_files_context(uploaded_files) -> str:
    if not uploaded_files:
        return ""

    chunks = [read_uploaded_file(f) for f in uploaded_files]
    return "\n\n".join(chunks)


def make_research_prompt(messages: List[Dict[str, str]], files_context: str, uploaded_files) -> str:
    user_question = ""

    for m in reversed(messages):
        if m.get("role") == "user":
            user_question = str(m.get("content", ""))
            break

    research_prompt = [
        build_system_prompt(),
        "",
        "Atlik gilią analizę lietuviškai. Pateik aiškią struktūrą, svarbiausias išvadas ir šaltinius, jei agentas juos turi.",
        "",
        "Vartotojo klausimas:",
        user_question,
    ]

    if files_context:
        research_prompt.extend(["", "Prisegtų failų tekstinis turinys:", files_context])

    image_names = [f.name for f in uploaded_files or [] if is_image_file(f)]

    if image_names:
        research_prompt.extend(
            [
                "",
                "Pastaba: ši pirmoji Gilios analizės versija vaizdų tiesiogiai Deep Research agentui dar neperduoda.",
                "Prisegti vaizdai: " + ", ".join(image_names),
            ]
        )

    return "\n".join(research_prompt).strip()


def extract_interaction_text(interaction) -> str:
    output_text = getattr(interaction, "output_text", "") or ""

    if output_text:
        return output_text

    outputs = getattr(interaction, "outputs", None) or []

    if outputs:
        last_output = outputs[-1]
        text = getattr(last_output, "text", "") or ""

        if text:
            return text

    steps = getattr(interaction, "steps", None) or []

    for step in reversed(steps):
        text = getattr(step, "text", "") or ""

        if text:
            return text

        output = getattr(step, "output", None)
        text = getattr(output, "text", "") or ""

        if text:
            return text

    return ""


def stream_gemini_deep_research(files_context: str, uploaded_files=None):
    api_key = get_api_key("Google Gemini")

    if not api_key:
        yield "🔴 Google Gemini API raktas neįvestas Streamlit Secrets."
        return

    if genai is None:
        yield "🔴 Deep Research reikia google-genai bibliotekos."
        return

    if st.session_state.provider != "Google Gemini":
        yield "🔴 Gili analizė šiuo metu veikia tik pasirinkus Google Gemini teikėją."
        return

    agent = DEEP_RESEARCH_AGENT_FAST
    mode_label = "Gili analizė"

    client = genai.Client(api_key=api_key)
    research_input = make_research_prompt(st.session_state.messages, files_context, uploaded_files or [])

    try:
        interaction = client.interactions.create(
            input=research_input,
            agent=agent,
            background=True,
        )
    except Exception as exc:
        yield f"⚠️ Nepavyko paleisti {mode_label} užduoties: {exc}"
        return

    interaction_id = getattr(interaction, "id", "") or ""

    yield (
        f"🔎 **{mode_label} pradėta.**\n\n"
        "Gemini Deep Research renka informaciją ir ruošia atsakymą. "
        "Tai gali trukti ilgiau nei paprastas pokalbis.\n\n"
    )

    for _ in range(DEEP_RESEARCH_MAX_POLLS):
        time.sleep(DEEP_RESEARCH_POLL_SECONDS)

        try:
            interaction = client.interactions.get(interaction_id)
        except Exception as exc:
            yield f"\n⚠️ Nepavyko patikrinti tyrimo būsenos: {exc}"
            return

        status = str(getattr(interaction, "status", "") or "").lower()

        if status == "completed":
            result = extract_interaction_text(interaction)
            yield "\n✅ **Tyrimas baigtas.**\n\n"
            yield result or "Gautas tuščias Deep Research atsakymas."
            return

        if status == "failed":
            error = getattr(interaction, "error", "") or "Nežinoma klaida."
            yield f"\n⚠️ Tyrimas nepavyko: {error}"
            return

        yield "⏳ "

    yield (
        "\n⚠️ Tyrimas dar nebaigtas per numatytą laiką. "
        "Pabandyk trumpesnį klausimą arba naudok Mąstymas režimą, kuris yra greitesnis ir neveikia per Deep Research."
    )


def make_google_config():
    system_instruction = build_system_prompt()

    if genai_types is not None:
        try:
            tools = None

            if st.session_state.web_search_enabled:
                tools = [genai_types.Tool(google_search=genai_types.GoogleSearch())]

            return genai_types.GenerateContentConfig(
                temperature=st.session_state.temperature,
                system_instruction=system_instruction,
                tools=tools,
            )
        except Exception:
            pass

    return {
        "temperature": st.session_state.temperature,
        "system_instruction": system_instruction,
    }


def answer_google(model: str, messages: List[Dict[str, str]], files_context: str, uploaded_files=None) -> str:
    api_key = get_api_key("Google Gemini")

    if not api_key:
        return "🔴 Google Gemini API raktas neįvestas Streamlit Secrets."

    client = genai.Client(api_key=api_key)
    contents = make_google_contents(messages, files_context, uploaded_files or [])

    config = make_google_config()

    response = client.models.generate_content(
        model=model,
        contents=contents,
        config=config,
    )

    return getattr(response, "text", "") or "Gautas tuščias atsakymas."


def stream_google(model: str, messages: List[Dict[str, str]], files_context: str, uploaded_files=None):
    api_key = get_api_key("Google Gemini")

    if not api_key:
        yield "🔴 Google Gemini API raktas neįvestas Streamlit Secrets."
        return

    client = genai.Client(api_key=api_key)
    contents = make_google_contents(messages, files_context, uploaded_files or [])

    config = make_google_config()

    try:
        stream = client.models.generate_content_stream(
            model=model,
            contents=contents,
            config=config,
        )

        for chunk in stream:
            text = getattr(chunk, "text", "") or ""
            if text:
                yield text
    except Exception:
        yield answer_google(model, messages, files_context, uploaded_files)


def answer_openai(model: str, messages: List[Dict[str, str]], files_context: str, uploaded_files=None) -> str:
    api_key = get_api_key("OpenAI")

    if not api_key:
        return "🔴 OpenAI API raktas neįvestas Streamlit Secrets."

    client = OpenAI(api_key=api_key)
    input_messages = make_openai_input_messages(messages, files_context, uploaded_files or [])

    kwargs = {
        "model": model,
        "input": input_messages,
    }

    if not model.startswith(("gpt-5", "o1", "o3", "o4")):
        kwargs["temperature"] = st.session_state.temperature

    if st.session_state.web_search_enabled:
        kwargs["tools"] = [{"type": "web_search"}]

    try:
        response = client.responses.create(**kwargs)
    except Exception as exc:
        msg = str(exc)

        if "temperature" in msg and "Unsupported parameter" in msg:
            kwargs.pop("temperature", None)
            response = client.responses.create(**kwargs)
        elif "web_search" in msg or "tool" in msg.lower():
            kwargs.pop("tools", None)
            response = client.responses.create(**kwargs)
        else:
            raise

    return getattr(response, "output_text", "") or "Gautas tuščias atsakymas."


def stream_openai(model: str, messages: List[Dict[str, str]], files_context: str, uploaded_files=None):
    api_key = get_api_key("OpenAI")

    if not api_key:
        yield "🔴 OpenAI API raktas neįvestas Streamlit Secrets."
        return

    client = OpenAI(api_key=api_key)
    input_messages = make_openai_input_messages(messages, files_context, uploaded_files or [])

    kwargs = {
        "model": model,
        "input": input_messages,
        "stream": True,
    }

    if not model.startswith(("gpt-5", "o1", "o3", "o4")):
        kwargs["temperature"] = st.session_state.temperature

    if st.session_state.web_search_enabled:
        kwargs["tools"] = [{"type": "web_search"}]

    try:
        stream = client.responses.create(**kwargs)

        for event in stream:
            event_type = getattr(event, "type", "")

            if event_type == "response.output_text.delta":
                delta = getattr(event, "delta", "") or ""
                if delta:
                    yield delta

    except Exception as exc:
        msg = str(exc)

        try:
            if "temperature" in msg and "Unsupported parameter" in msg:
                kwargs.pop("temperature", None)
                stream = client.responses.create(**kwargs)

                for event in stream:
                    event_type = getattr(event, "type", "")

                    if event_type == "response.output_text.delta":
                        delta = getattr(event, "delta", "") or ""
                        if delta:
                            yield delta

            elif "web_search" in msg or "tool" in msg.lower():
                kwargs.pop("tools", None)
                stream = client.responses.create(**kwargs)

                for event in stream:
                    event_type = getattr(event, "type", "")

                    if event_type == "response.output_text.delta":
                        delta = getattr(event, "delta", "") or ""
                        if delta:
                            yield delta
            else:
                raise
        except Exception:
            kwargs.pop("stream", None)
            yield answer_openai(model, messages, files_context, uploaded_files)


def stream_answer(uploaded_files):
    files_context = make_files_context(uploaded_files)
    provider = st.session_state.provider
    model = st.session_state.model

    try:
        if st.session_state.mode == "deep":
            yield from stream_gemini_deep_research(files_context, uploaded_files)
            return

        if provider == "Google Gemini":
            yield from stream_google(model, st.session_state.messages, files_context, uploaded_files)
            return

        if provider == "OpenAI":
            yield from stream_openai(model, st.session_state.messages, files_context, uploaded_files)
            return

        yield "Nežinomas teikėjas."
    except Exception as exc:
        yield f"⚠️ Klaida: {exc}"


def render_brand() -> None:
    logo_path = "chatmd_logo.png"
    logo_html = ""

    try:
        with open(logo_path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode("utf-8")
        logo_html = f'<img class="chat-logo" src="data:image/png;base64,{b64}" />'
    except Exception:
        logo_html = ""

    st.sidebar.markdown(
        f"""
        <div class="chat-title">
            {logo_html}
            <div>
                <div class="brand-main">ChatMD</div>
                <div class="brand-sub">AI pokalbiai ir agentai</div>
            </div>
        </div>
        <div class="version-pill">{APP_VERSION}</div>
        """,
        unsafe_allow_html=True,
    )


def render_agents_panel() -> None:
    st.session_state.agents = load_agents_from_db()
    agents = st.session_state.agents

    st.sidebar.subheader("Sukurti agentai")

    if not agents:
        st.sidebar.caption("Agentų dar nėra.")

    agent_names = list(agents.keys())

    if st.session_state.agent_name not in agent_names and agent_names:
        st.session_state.agent_name = agent_names[0]

    if agent_names:
        selected_agent = st.sidebar.radio(
            "Pasirink agentą",
            agent_names,
            index=agent_names.index(st.session_state.agent_name)
            if st.session_state.agent_name in agent_names
            else 0,
            format_func=lambda name: f"{agents[name].get('icon', '🤖')} {name}",
            label_visibility="collapsed",
            key="agent_selector_radio",
        )

        if selected_agent != st.session_state.agent_name:
            st.session_state.agent_name = selected_agent
            reset_chat()
            st.rerun()

        active = agents.get(st.session_state.agent_name, {})
        st.sidebar.caption(
            f"Aktyvus agentas: {active.get('icon', '🤖')} {st.session_state.agent_name}"
        )

    with st.sidebar.expander("⚙️ Agentų valdymas", expanded=False):
        st.caption("Čia gali kurti, redaguoti, trinti, importuoti ir eksportuoti agentus.")

        with st.expander("➕ Sukurti naują agentą", expanded=False):
            new_icon = st.text_input("Ikona", value="🤖", key="new_agent_icon")
            new_name = st.text_input("Pavadinimas", value="", key="new_agent_name")
            new_prompt = st.text_area(
                "Taisyklės / sisteminė instrukcija",
                value="",
                height=160,
                key="new_agent_prompt",
                placeholder="Pvz.: Tu esi teisinis asistentas. Atsakyk aiškiai, struktūruotai, lietuviškai...",
            )
            new_examples = st.text_area(
                "Pavyzdžiai / papildomos žinios",
                value="",
                height=120,
                key="new_agent_examples",
                placeholder="Čia gali įrašyti atsakymų pavyzdžius, toną, darbo tvarką arba specifines žinias.",
            )

            if st.button("Išsaugoti naują agentą", use_container_width=True):
                try:
                    save_agent_to_db(new_name, new_icon, new_prompt, new_examples)
                    st.session_state.agent_name = new_name.strip()
                    st.session_state.agents = load_agents_from_db()
                    st.success("Agentas sukurtas.")
                    st.rerun()
                except Exception as exc:
                    st.error(f"Nepavyko sukurti agento: {exc}")

        if agent_names and st.session_state.agent_name in agents:
            current = agents[st.session_state.agent_name]

            with st.expander("✏️ Redaguoti aktyvų agentą", expanded=False):
                edit_icon = st.text_input(
                    "Ikona",
                    value=current.get("icon", "🤖"),
                    key=f"edit_agent_icon_{st.session_state.agent_name}",
                )
                edit_name = st.text_input(
                    "Pavadinimas",
                    value=st.session_state.agent_name,
                    key=f"edit_agent_name_{st.session_state.agent_name}",
                )
                edit_prompt = st.text_area(
                    "Taisyklės / sisteminė instrukcija",
                    value=current.get("system_prompt", ""),
                    height=180,
                    key=f"edit_agent_prompt_{st.session_state.agent_name}",
                )
                edit_examples = st.text_area(
                    "Pavyzdžiai / papildomos žinios",
                    value=current.get("examples", ""),
                    height=140,
                    key=f"edit_agent_examples_{st.session_state.agent_name}",
                )

                if st.button("Išsaugoti pakeitimus", use_container_width=True):
                    old_name = st.session_state.agent_name
                    try:
                        save_agent_to_db(
                            edit_name,
                            edit_icon,
                            edit_prompt,
                            edit_examples,
                            old_name=old_name,
                        )
                        st.session_state.agent_name = edit_name.strip()
                        st.session_state.agents = load_agents_from_db()
                        st.success("Agentas atnaujintas.")
                        st.rerun()
                    except Exception as exc:
                        st.error(f"Nepavyko atnaujinti agento: {exc}")

            with st.expander("🗑 Ištrinti aktyvų agentą", expanded=False):
                st.warning("Ištrynus agentą, jo pokalbių istorija liks, bet pats agentas bus pašalintas.")
                confirm_delete = st.checkbox(
                    "Patvirtinu, kad noriu ištrinti šį agentą",
                    key=f"confirm_delete_agent_{st.session_state.agent_name}",
                )

                if st.button("Ištrinti agentą", use_container_width=True):
                    if not confirm_delete:
                        st.error("Pirma pažymėk patvirtinimo varnelę.")
                    else:
                        try:
                            delete_agent_from_db(st.session_state.agent_name)
                            st.session_state.agents = load_agents_from_db()
                            remaining = list(st.session_state.agents.keys())
                            st.session_state.agent_name = remaining[0] if remaining else "Bendras asistentas"
                            st.success("Agentas ištrintas.")
                            st.rerun()
                        except Exception as exc:
                            st.error(f"Nepavyko ištrinti agento: {exc}")

        with st.expander("💾 Importas / eksportas", expanded=False):
            st.caption("Eksportas išsaugo tik agentus. API raktai ir slaptažodžiai neeksportuojami.")

            export_json = make_agents_export_json()

            st.download_button(
                "Eksportuoti agentus į JSON",
                data=export_json,
                file_name="chatmd_agents.json",
                mime="application/json",
                use_container_width=True,
            )

            uploaded_agents_json = st.file_uploader(
                "Importuoti agentus iš JSON",
                type=["json"],
                key="agents_import_json",
            )

            if uploaded_agents_json is not None:
                if st.button("Importuoti agentus", use_container_width=True):
                    try:
                        imported_count = import_agents_from_json_bytes(uploaded_agents_json.getvalue())
                        st.session_state.agents = load_agents_from_db()
                        agent_names_after_import = list(st.session_state.agents.keys())

                        if agent_names_after_import:
                            st.session_state.agent_name = agent_names_after_import[0]

                        st.success(f"Importuota agentų: {imported_count}")
                        st.rerun()
                    except Exception as exc:
                        st.error(f"Nepavyko importuoti agentų: {exc}")


def render_clipboard_paste_panel() -> None:
    pasted_count = len(st.session_state.get("pasted_images", []))

    if paste_image_button is None:
        return

    st.markdown(
        """
        <style>
        .paste-toolbar-spacer {
            height: clamp(310px, 61vh, 700px);
        }

        @media (max-height: 760px) {
            .paste-toolbar-spacer {
                height: clamp(190px, 46vh, 430px);
            }
        }
        </style>
        <div class="paste-toolbar-spacer"></div>
        """,
        unsafe_allow_html=True,
    )

    spacer, button_col, status_col, rest = st.columns([0.24, 0.16, 0.14, 0.46])

    with button_col:
        paste_result = paste_image_button(
            label="📋 Įklijuoti screenshot",
            text_color="#0f172a",
            background_color="#e8f1ff",
            hover_background_color="#dbeafe",
            key="clipboard_paste_button",
            errors="ignore",
        )

    pasted_image = getattr(paste_result, "image_data", None)

    if pasted_image is not None:
        added = add_pasted_image_to_state(pasted_image)
        if added:
            pasted_count = len(st.session_state.get("pasted_images", []))

    with status_col:
        if pasted_count:
            st.markdown(
                f'<span class="status-pill">🟢 {pasted_count} screenshot</span>',
                unsafe_allow_html=True,
            )
        else:
            st.markdown(
                '<span class="status-pill">⚪ Screenshot nėra</span>',
                unsafe_allow_html=True,
            )


def render_sidebar() -> None:
    render_brand()

    if st.sidebar.button("＋ Naujas pokalbis", use_container_width=True):
        reset_chat()
        st.rerun()

    st.sidebar.divider()

    st.sidebar.subheader("Režimas")
    mode = st.sidebar.radio(
        "Pasirink režimą",
        ["chat", "thinking", "agent", "deep"],
        format_func=lambda x: {
            "chat": "💬 Pokalbis",
            "thinking": "🧠 Mąstymas",
            "agent": "🤖 Agentai",
            "deep": "🔎 Gili analizė",
        }.get(x, x),
        horizontal=False,
        label_visibility="collapsed",
    )

    if mode != st.session_state.mode:
        reset_chat()
        st.session_state.mode = mode
        st.rerun()

    if st.session_state.mode == "thinking":
        st.sidebar.info("Mąstymas atsako giliau, bet veikia greitai ir nenaudoja Deep Research.")

    if st.session_state.mode == "agent":
        render_agents_panel()

    if st.session_state.mode == "deep":
        st.sidebar.info("Gili analizė naudoja Gemini Deep Research. Ji gali trukti kelias minutes. Pasirink Google Gemini teikėją.")

    st.sidebar.divider()

    st.sidebar.subheader("AI nustatymai")

    provider = st.sidebar.selectbox(
        "Teikėjas",
        ["Google Gemini", "OpenAI"],
        index=["Google Gemini", "OpenAI"].index(st.session_state.provider),
    )

    if provider != st.session_state.provider:
        st.session_state.provider = provider
        models, _ = list_models(provider)
        st.session_state.model = models[0] if models else ""
        reset_chat()
        st.rerun()

    key = get_api_key(st.session_state.provider)
    st.sidebar.markdown("🟢 API raktas aktyvus" if key else "🔴 API raktas neįvestas")

    if st.sidebar.button("Atnaujinti modelių sąrašą", use_container_width=True):
        models, source = list_models(st.session_state.provider)
        st.session_state.models_cache[st.session_state.provider] = models

        if st.session_state.model not in models:
            st.session_state.model = models[0] if models else ""

        st.sidebar.info(f"Modeliai: {source}")

    models = st.session_state.models_cache.get(st.session_state.provider)

    if not models:
        models, _ = list_models(st.session_state.provider)
        st.session_state.models_cache[st.session_state.provider] = models

    if st.session_state.model not in models and models:
        st.session_state.model = models[0]

    st.session_state.model = st.sidebar.selectbox(
        "Modelis",
        models,
        index=models.index(st.session_state.model) if st.session_state.model in models else 0,
    )

    st.session_state.web_search_enabled = st.sidebar.toggle(
        "Interneto paieška",
        value=st.session_state.web_search_enabled,
    )

    st.session_state.temperature = st.sidebar.slider(
        "Kūrybiškumas",
        0.0,
        2.0,
        float(st.session_state.temperature),
        0.1,
    )

    st.sidebar.divider()
    st.sidebar.subheader("Pokalbių istorija")

    saved_chats = list_saved_chats(limit=30)

    if not saved_chats:
        st.sidebar.caption("Istorija atsiras po pirmo pokalbio.")

    for chat in saved_chats:
        title = str(chat.get("title", "Pokalbis") or "Pokalbis").strip()
        chat_mode = chat.get("mode", "chat")

        mode_icon = {
            "chat": "💬",
            "thinking": "🧠",
            "agent": "🤖",
            "deep": "🔎",
        }.get(chat_mode, "💬")

        if len(title) > 30:
            title = title[:27].rstrip() + "..."

        label = f"{mode_icon} {title}"

        col1, col2 = st.sidebar.columns([0.82, 0.18])

        with col1:
            if st.button(label, key=f"hist_{chat['id']}", use_container_width=True):
                save_current_chat_to_db()
                loaded_chat = load_chat_from_db(chat["id"])

                if loaded_chat:
                    st.session_state.current_chat_id = loaded_chat["id"]
                    st.session_state.messages = list(loaded_chat.get("messages", []))
                    st.session_state.mode = loaded_chat.get("mode", "chat")
                    st.session_state.provider = loaded_chat.get("provider", st.session_state.provider)

                    loaded_agent_name = loaded_chat.get("agent_name", "")
                    if loaded_agent_name:
                        st.session_state.agent_name = loaded_agent_name

                    saved_model = loaded_chat.get("model", "")

                    if saved_model:
                        st.session_state.model = saved_model

                    st.rerun()

        with col2:
            if st.button("🗑", key=f"del_{chat['id']}"):
                delete_chat_from_db(chat["id"])

                if st.session_state.get("current_chat_id") == chat["id"]:
                    st.session_state.messages = []
                    st.session_state.current_chat_id = None

                st.rerun()


def parse_chat_input(chat_value):
    if not chat_value:
        return "", []

    if isinstance(chat_value, str):
        return chat_value.strip(), []

    text = getattr(chat_value, "text", "")
    files = getattr(chat_value, "files", [])

    if not text and isinstance(chat_value, dict):
        text = chat_value.get("text", "")
        files = chat_value.get("files", [])

    return str(text or "").strip(), list(files or [])


def render_user_message(content: str) -> None:
    safe_content = html.escape(content).replace("\n", "<br>")

    st.markdown(
        f"""
        <div class="msg-row user">
            <div class="msg-bubble user">{safe_content}</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_chat_message(role: str, content: str) -> None:
    if role == "user":
        render_user_message(content)
    else:
        st.markdown(content)


def render_thinking(placeholder) -> None:
    placeholder.markdown(
        """
        <div class="thinking-wrap">
            <div class="thinking-card">
                <div class="thinking-spinner"></div>
                <div>
                    <div class="thinking-text">
                        Galvoju<span class="thinking-dots"><span>.</span><span>.</span><span>.</span></span>
                    </div>
                    <div class="thinking-time">Ruošiu atsakymą</div>
                </div>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def main() -> None:
    st.set_page_config(
        page_title=APP_NAME,
        page_icon="💬",
        layout="wide",
        initial_sidebar_state="expanded",
    )

    init_history_db()
    init_state()

    if not check_password():
        return

    # Atstatome paskutinio pokalbio modelį ir teikėją tik vieną kartą po programos paleidimo.
    # Taip šoniniame meniu ranka pakeistas modelis nebėra perrašomas kiekvieno Streamlit rerun metu.
    if (
        not st.session_state.last_chat_settings_loaded
        and not st.session_state.messages
        and not st.session_state.current_chat_id
    ):
        saved_chats = list_saved_chats(limit=1)
        if saved_chats:
            last_chat = load_chat_from_db(saved_chats[0]["id"])
            if last_chat:
                saved_model = last_chat.get("model", "")
                saved_provider = last_chat.get("provider", "")
                saved_mode = last_chat.get("mode", "chat")
                saved_agent = last_chat.get("agent_name", "")

                if saved_provider:
                    st.session_state.provider = saved_provider
                if saved_model:
                    st.session_state.model = saved_model
                if saved_mode:
                    st.session_state.mode = saved_mode
                if saved_agent:
                    st.session_state.agent_name = saved_agent

        st.session_state.last_chat_settings_loaded = True

    render_css()
    render_sidebar()

    mode_label = {
        "chat": "Pokalbis",
        "thinking": "Mąstymas",
        "agent": f"Agentai: {st.session_state.agent_name}",
        "deep": "Gili analizė",
    }.get(st.session_state.mode, "Pokalbis")

    st.markdown("### ChatMD")
    st.markdown(
        f"""
        <span class="version-pill">{APP_VERSION}</span>
        <br>
        <span class="status-pill">{st.session_state.provider}</span>
        <span class="status-pill">{st.session_state.model}</span>
        <span class="status-pill">{mode_label}</span>
        <span class="status-pill">Internetas: {"ON" if st.session_state.web_search_enabled else "OFF"}</span>
        """,
        unsafe_allow_html=True,
    )

    st.markdown('<div class="chat-area">', unsafe_allow_html=True)

    for m in st.session_state.messages:
        render_chat_message(m["role"], m["content"])

    st.markdown("</div>", unsafe_allow_html=True)

    render_clipboard_paste_panel()

    chat_value = st.chat_input(
        "Rašykite žinutę...",
        accept_file="multiple",
        file_type=[
            "png",
            "jpg",
            "jpeg",
            "webp",
            "gif",
            "pdf",
            "docx",
            "xlsx",
            "pptx",
            "txt",
            "csv",
            "md",
            "py",
            "json",
        ],
    )

    prompt, uploaded_files = parse_chat_input(chat_value)

    if prompt or uploaded_files:
        uploaded_files = list(uploaded_files or []) + get_pasted_image_files()

        if not prompt and uploaded_files:
            prompt = "Peržiūrėk prisegtus failus."

        visible_prompt = prompt

        if uploaded_files:
            file_names = ", ".join([f.name for f in uploaded_files])
            visible_prompt = f"{prompt}\n\n📎 Prisegti failai: {file_names}"

        st.session_state.messages.append({"role": "user", "content": visible_prompt})

        render_user_message(visible_prompt)

        assistant_placeholder = st.empty()
        render_thinking(assistant_placeholder)

        streamed_answer = ""

        for chunk in stream_answer(uploaded_files):
            streamed_answer += chunk
            assistant_placeholder.markdown(streamed_answer + "▌")

        final_answer = streamed_answer.strip() or "Gautas tuščias atsakymas."
        assistant_placeholder.markdown(final_answer)

        st.session_state.messages.append({"role": "assistant", "content": final_answer})
        save_current_chat_to_db()
        clear_pasted_images()


if __name__ == "__main__":
    main()
